import os
import logging
import tempfile
import asyncio
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from fpdf import FPDF
import aiohttp
from io import BytesIO
from PIL import Image

logger = logging.getLogger(__name__)

class PresentationGenerator:
    def __init__(self):
        api_key = os.getenv('OPENAI_API_KEY')
        if not api_key:
            raise ValueError("OPENAI_API_KEY not found in environment variables!")
        if not api_key.startswith('sk-'):
            raise ValueError("OPENAI_API_KEY appears to be invalid!")
        self.client = OpenAI(api_key=api_key)
        
        self.background_images = {
            'asosiy': 'slayd_fon/asosiy_sahifa.png',
            'reja': 'slayd_fon/orta_sahifa.png',
            'content_1': 'slayd_fon/2.png',
            'content_2': 'slayd_fon/3.png',
            'content_3': 'slayd_fon/4.png',
            'oxirgi': 'slayd_fon/oxirgi_sahifa.png'
        }
    
    async def generate_presentation(self, topic: str, num_slides: int, plan: str):
        logger.info(f"Generating presentation: {topic}, {num_slides} slides, {plan} plan")
        
        slides_content = await self.generate_slides_content(topic, num_slides)
        
        ppt_path = await self.create_ppt(topic, slides_content)
        
        files = [ppt_path]
        
        if plan == 'smart':
            pdf_path = await self.create_pdf(topic, slides_content)
            files.append(pdf_path)
        
        return files
    
    async def generate_slides_content(self, topic: str, num_slides: int):
        logger.info(f"Generating content for {num_slides} slides")
        
        content_slides = num_slides - 3
        
        prompt = f"""
Taqdimot mavzusi: {topic}

Quyidagi struktura bo'yicha {num_slides} ta slayd uchun kontent yarating:

1. KIRISH (1 slayd) - rasm bor, bullet points
2. REJA (1 slayd) - faqat reja
3. ASOSIY QISM ({content_slides} ta slayd):
   - Dastlabki 2 ta slayd: rasm bor, 4-5 ta bullet point
   - Qolgan slaydlar: rasm YO'Q, PARAGRAF shaklida batafsil matn
4. XULOSA (1 slayd) - rasm yo'q, PARAGRAF shaklida

MUHIM: Har bir slayd uchun SLIDE belgisini yozing!

Format (ANIQ SHU FORMATDA YOZING):

SLIDE kirish
TITLE: Kirish
CONTENT:
- [Nuqta 1]
- [Nuqta 2]
- [Nuqta 3]
- [Nuqta 4]
IMAGE_PROMPT: [Professional illustration for introduction]

SLIDE reja
SECTION_1: [1-bo'lim nomi]
SECTION_2: [2-bo'lim nomi]
SECTION_3: [3-bo'lim nomi]

SLIDE 1
TITLE: [Sarlavha]
CONTENT:
- [Nuqta 1]
- [Nuqta 2]
- [Nuqta 3]
- [Nuqta 4]
IMAGE_PROMPT: [Image description]

SLIDE 2
TITLE: [Sarlavha]
CONTENT:
- [Nuqta 1]
- [Nuqta 2]
- [Nuqta 3]
- [Nuqta 4]
IMAGE_PROMPT: [Image description]

SLIDE 3
TITLE: [Sarlavha]
CONTENT:
[Mavzuga oid batafsil va professional matn yozing. Kamida 150-200 so'z. Bir necha paragraf bo'lsin. Har bir paragraf alohida qatordan boshlansin. Bullet point ishlatmasdan, oddiy matn formatida yozing.]

SLIDE 4
TITLE: [Sarlavha]
CONTENT:
[Mavzuga oid batafsil va professional matn yozing. Kamida 150-200 so'z. Bir necha paragraf bo'lsin. Har bir paragraf alohida qatordan boshlansin. Bullet point ishlatmasdan, oddiy matn formatida yozing.]

[Qolgan slaydlar ham xuddi shunday formatda...]

SLIDE xulosa
TITLE: Xulosa
CONTENT:
[Yakuniy xulosalar va fikrlarni batafsil yozing. Kamida 100-150 so'z. Bir necha paragraf bo'lsin. Bullet point ishlatmasdan, oddiy matn formatida yozing.]

JAMI {num_slides} TA SLAYD BO'LISHI KERAK!
"""
        
        try:
            response = self.client.chat.completions.create(
                model="gpt-4.1",
                messages=[
                    {"role": "system", "content": "Siz professional taqdimot yaratuvchi AI assistentsiz. O'zbek tilida yozing, lekin rasm tavsiflari ingliz tilida bo'lsin. MUHIM: [Kvadrat qavs ichidagi] ko'rsatmalarni YOZMASDAN, ularning o'rniga HAQIQIY KONTENT yozing!"},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.7,
                max_tokens=4000
            )
            
            content = response.choices[0].message.content
            slides = self.parse_slides_content(content)
            
            ai_image_count = 0
            for slide in slides:
                if slide.get('type') != 'reja' and 'image_prompt' in slide and slide['image_prompt'] and ai_image_count < 3:
                    slide['image_url'] = await self.generate_image(slide['image_prompt'])
                    ai_image_count += 1
                else:
                    slide['image_url'] = None
            
            return slides
            
        except Exception as e:
            logger.error(f"Error generating slides content: {e}")
            raise
    
    async def generate_image(self, prompt: str):
        logger.info(f"Generating image for: {prompt}")
        
        try:
            response = self.client.images.generate(
                model="dall-e-3",
                prompt=f"Professional presentation slide image: {prompt}. Clean, modern, business style.",
                size="1024x1024",
                quality="standard",
                n=1
            )
            
            image_url = response.data[0].url
            logger.info(f"Image generated: {image_url}")
            return image_url
            
        except Exception as e:
            logger.error(f"Error generating image: {e}")
            return None
    
    def parse_slides_content(self, content: str):
        slides = []
        current_slide = {}
        
        lines = content.strip().split('\n')
        
        for line in lines:
            line = line.strip()
            
            if line.startswith('SLIDE '):
                if current_slide and current_slide.get('type'):
                    slides.append(current_slide)
                
                slide_type = line.replace('SLIDE ', '').strip()
                current_slide = {'type': slide_type, 'title': '', 'content': [], 'image_prompt': ''}
                
                if slide_type == 'reja':
                    current_slide['sections'] = []
            
            elif line.startswith('TITLE:'):
                if not current_slide:
                    current_slide = {'type': 'content', 'title': '', 'content': [], 'image_prompt': ''}
                current_slide['title'] = line.replace('TITLE:', '').strip()
            
            elif line.startswith('SECTION_'):
                if not current_slide:
                    current_slide = {'type': 'reja', 'title': '', 'content': [], 'sections': []}
                if current_slide.get('type') == 'reja':
                    current_slide['sections'].append(line.split(':', 1)[1].strip())
            
            elif line.startswith('CONTENT:'):
                continue
            
            elif line.startswith('-'):
                if not current_slide:
                    current_slide = {'type': 'content', 'title': '', 'content': [], 'image_prompt': ''}
                current_slide['content'].append(line[1:].strip())
            
            elif line.startswith('IMAGE_PROMPT:'):
                if not current_slide:
                    current_slide = {'type': 'content', 'title': '', 'content': [], 'image_prompt': ''}
                current_slide['image_prompt'] = line.replace('IMAGE_PROMPT:', '').strip()
            
            elif line and current_slide.get('type') and not line.startswith('[') and not line.startswith('('):
                if not line.startswith('SLIDE') and not line.startswith('TITLE:') and not line.startswith('SECTION_') and not line.startswith('CONTENT:') and not line.startswith('IMAGE_PROMPT:') and not line.startswith('JAMI'):
                    if 'content' in current_slide:
                        current_slide['content'].append(line)
        
        if current_slide and current_slide.get('type'):
            slides.append(current_slide)
        
        logger.info(f"Parsed {len(slides)} slides from GPT response")
        
        if len(slides) == 0:
            logger.error(f"Failed to parse slides. GPT response:\n{content}")
            raise ValueError("Could not parse slides from GPT response")
        
        return slides
    
    async def create_ppt(self, topic: str, slides_content: list):
        logger.info("Creating PowerPoint presentation")
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        blank_layout = prs.slide_layouts[6]
        
        slide = prs.slides.add_slide(blank_layout)
        if os.path.exists(self.background_images['asosiy']):
            slide.shapes.add_picture(
                self.background_images['asosiy'],
                0, 0,
                width=prs.slide_width,
                height=prs.slide_height
            )
        
        title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = topic
        title_frame.paragraphs[0].font.size = Pt(48)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "@preuz_bot"
        subtitle_frame.paragraphs[0].font.size = Pt(20)
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        for idx, slide_data in enumerate(slides_content):
            slide = prs.slides.add_slide(blank_layout)
            
            if slide_data.get('type') == 'reja':
                if os.path.exists(self.background_images['reja']):
                    slide.shapes.add_picture(
                        self.background_images['reja'],
                        0, 0,
                        width=prs.slide_width,
                        height=prs.slide_height
                    )
                
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
                title_frame = title_box.text_frame
                title_frame.text = "REJA"
                title_frame.paragraphs[0].font.size = Pt(36)
                title_frame.paragraphs[0].font.bold = True
                title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
                sections = slide_data.get('sections', [])
                box_width = Inches(2.5)
                box_height = Inches(2)
                start_x = Inches(1)
                start_y = Inches(3)
                spacing = Inches(0.5)
                
                for i, section in enumerate(sections[:3]):
                    x = start_x + (i * (box_width + spacing))
                    box = slide.shapes.add_textbox(x, start_y, box_width, box_height)
                    text_frame = box.text_frame
                    text_frame.word_wrap = True
                    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                    
                    p = text_frame.paragraphs[0]
                    p.text = f"{i+1}. {section}"
                    p.font.size = Pt(16)
                    p.font.bold = True
                    p.alignment = PP_ALIGN.CENTER
            
            elif slide_data.get('type') == 'xulosa':
                if os.path.exists(self.background_images['oxirgi']):
                    slide.shapes.add_picture(
                        self.background_images['oxirgi'],
                        0, 0,
                        width=prs.slide_width,
                        height=prs.slide_height
                    )
                
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
                title_frame = title_box.text_frame
                title_frame.text = slide_data.get('title', 'Xulosa')
                title_frame.paragraphs[0].font.size = Pt(36)
                title_frame.paragraphs[0].font.bold = True
                
                content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
                content_frame = content_box.text_frame
                content_frame.word_wrap = True
                
                for i, point in enumerate(slide_data.get('content', [])):
                    if i > 0:
                        content_frame.add_paragraph()
                    p = content_frame.paragraphs[i]
                    p.text = point
                    p.font.size = Pt(18)
                    p.space_after = Pt(12)
            
            else:
                bg_key = f'content_{(idx % 3) + 1}'
                if os.path.exists(self.background_images.get(bg_key, '')):
                    slide.shapes.add_picture(
                        self.background_images[bg_key],
                        0, 0,
                        width=prs.slide_width,
                        height=prs.slide_height
                    )
                
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
                title_frame = title_box.text_frame
                title_frame.text = slide_data.get('title', '')
                title_frame.paragraphs[0].font.size = Pt(32)
                title_frame.paragraphs[0].font.bold = True
                
                has_image = bool(slide_data.get('image_url'))
                
                if has_image:
                    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(5), Inches(4.5))
                    
                    try:
                        async with aiohttp.ClientSession() as session:
                            async with session.get(slide_data['image_url']) as response:
                                image_data = await response.read()
                                image_stream = BytesIO(image_data)
                        
                        slide.shapes.add_picture(
                            image_stream,
                            Inches(6),
                            Inches(2),
                            width=Inches(3.5)
                        )
                    except Exception as e:
                        logger.error(f"Error adding image: {e}")
                        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(4.5))
                        has_image = False
                else:
                    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(4.5))
                
                content_frame = content_box.text_frame
                content_frame.word_wrap = True
                
                for i, point in enumerate(slide_data.get('content', [])):
                    if i > 0:
                        content_frame.add_paragraph()
                    p = content_frame.paragraphs[i]
                    
                    if has_image:
                        p.text = f"• {point}"
                    else:
                        p.text = point
                    
                    p.font.size = Pt(18)
                    p.space_after = Pt(12)
    
        # Presentations papkasini yaratish
        presentations_dir = "presentations"
        if not os.path.exists(presentations_dir):
            os.makedirs(presentations_dir)
        
        safe_topic = "".join(c if c.isalnum() or c in (' ', '_') else '_' for c in topic[:30])
        safe_topic = safe_topic.replace(' ', '_')
        
        # Doimiy fayl nomi
        import time
        timestamp = int(time.time())
        filename = os.path.join(presentations_dir, f"{safe_topic}_{timestamp}.pptx")
        
        prs.save(filename)
        logger.info(f"PowerPoint saved: {filename}")
        
        return filename
    
    async def create_pdf(self, topic: str, slides_content: list):
        logger.info("Creating PDF presentation")
        
        pdf = FPDF(orientation='L', unit='mm', format='A4')
        pdf.set_auto_page_break(auto=False)
        
        pdf.add_page()
        if os.path.exists(self.background_images['asosiy']):
            pdf.image(self.background_images['asosiy'], x=0, y=0, w=297, h=210)
        pdf.set_font('Arial', 'B', 32)
        pdf.ln(80)
        pdf.cell(0, 20, topic, align='C', ln=True)
        pdf.set_font('Arial', 'I', 14)
        pdf.cell(0, 10, '@preuz_bot', align='C')
        
        for idx, slide_data in enumerate(slides_content):
            pdf.add_page()
            
            if slide_data.get('type') == 'reja':
                if os.path.exists(self.background_images['reja']):
                    pdf.image(self.background_images['reja'], x=0, y=0, w=297, h=210)
                
                pdf.set_font('Arial', 'B', 28)
                pdf.cell(0, 30, 'REJA', align='C', ln=True)
                
                sections = slide_data.get('sections', [])
                pdf.set_font('Arial', 'B', 16)
                pdf.ln(20)
                for i, section in enumerate(sections[:3]):
                    pdf.cell(0, 15, f"{i+1}. {section}", align='C', ln=True)
            
            elif slide_data.get('type') == 'xulosa':
                if os.path.exists(self.background_images['oxirgi']):
                    pdf.image(self.background_images['oxirgi'], x=0, y=0, w=297, h=210)
                
                pdf.set_font('Arial', 'B', 24)
                pdf.cell(0, 20, slide_data.get('title', 'Xulosa'), ln=True)
                
                pdf.set_font('Arial', '', 12)
                pdf.ln(10)
                
                for point in slide_data.get('content', []):
                    pdf.multi_cell(0, 8, point)
                    pdf.ln(3)
            
            else:
                bg_key = f'content_{(idx % 3) + 1}'
                if os.path.exists(self.background_images.get(bg_key, '')):
                    pdf.image(self.background_images[bg_key], x=0, y=0, w=297, h=210)
                
                pdf.set_font('Arial', 'B', 20)
                pdf.cell(0, 20, slide_data.get('title', ''), ln=True)
                
                pdf.set_font('Arial', '', 12)
                pdf.ln(10)
                
                has_image = bool(slide_data.get('image_url'))
                
                for point in slide_data.get('content', []):
                    if has_image:
                        pdf.multi_cell(0, 8, f"  - {point}")
                    else:
                        pdf.multi_cell(0, 8, point)
                    pdf.ln(3)
                
                if has_image:
                    try:
                        async with aiohttp.ClientSession() as session:
                            async with session.get(slide_data['image_url']) as response:
                                image_data = await response.read()
                                image_stream = BytesIO(image_data)
                        
                        img = Image.open(image_stream)
                        
                        with tempfile.NamedTemporaryFile(delete=False, suffix='.jpg') as tmp_file:
                            img_path = tmp_file.name
                            img.save(img_path, 'JPEG')
                        
                        pdf.image(img_path, x=200, y=50, w=80)
                        
                        if os.path.exists(img_path):
                            os.remove(img_path)
                            
                    except Exception as e:
                        logger.error(f"Error adding image to PDF: {e}")
        
        # Presentations papkasini yaratish
        presentations_dir = "presentations"
        if not os.path.exists(presentations_dir):
            os.makedirs(presentations_dir)
        
        safe_topic = "".join(c if c.isalnum() or c in (' ', '_') else '_' for c in topic[:30])
        safe_topic = safe_topic.replace(' ', '_')
        
        # Doimiy fayl nomi
        import time
        timestamp = int(time.time())
        filename = os.path.join(presentations_dir, f"{safe_topic}_{timestamp}.pdf")
        
        pdf.output(filename)
        logger.info(f"PDF saved: {filename}")
        
        return filename

# Bot uchun wrapper funksiyalar
async def create_presentation_file(topic: str, num_slides: int, plan: str):
    """Bot uchun wrapper funksiya"""
    generator = PresentationGenerator()
    return await generator.generate_presentation(topic, num_slides, plan)

async def generate_presentation_content_with_gpt(topic: str, num_slides: int):
    """Bot uchun GPT kontent generator funksiya"""
    generator = PresentationGenerator()
    return await generator.generate_slides_content(topic, num_slides)